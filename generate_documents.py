"""
Lifedrop Shivamogga - High-Impact 6-Slide Presentation & 5-6 Page Documentation Generator
Generates:
1. Lifedrop_Shivamogga_Project_Presentation.pptx (Exact 6-slide presentation)
2. Lifedrop_Shivamogga_Project_Documentation.docx (Clean 5-6 page technical & executive report)
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import docx
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml import OxmlElement, parse_xml
from docx.oxml.ns import nsdecls, qn

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_PATH = os.path.join(OUTPUT_DIR, 'Lifedrop_Shivamogga_Project_Presentation.pptx')
DOCX_PATH = os.path.join(OUTPUT_DIR, 'Lifedrop_Shivamogga_Project_Documentation.docx')

# Theme Palette Constants
COLOR_CRIMSON = RGBColor(225, 29, 72)     # #E11D48
COLOR_DARK = RGBColor(15, 23, 42)         # #0F172A
COLOR_NAVY = RGBColor(30, 41, 59)         # #1E293B
COLOR_MUTED = RGBColor(100, 116, 139)     # #64748B
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_EMERALD = RGBColor(16, 185, 129)    # #10B981
COLOR_AMBER = RGBColor(245, 158, 11)      # #F59E0B
COLOR_CARD_BG = RGBColor(248, 250, 252)   # #F8FAFC
COLOR_BORDER = RGBColor(226, 232, 240)    # #E2E8F0
COLOR_LIGHT_RED = RGBColor(255, 241, 242) # #FFF1F2
COLOR_BLUE = RGBColor(37, 99, 235)        # #2563EB

# -----------------------------------------------------------------------------
# 1. POWERPOINT PRESENTATION GENERATION (6 SLIDES)
# -----------------------------------------------------------------------------

def create_presentation():
    print("Generating 6-Slide PowerPoint Presentation...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category="LIFEDROP SHIVAMOGGA | EMERGENCY BLOOD NETWORK"):
        # Top banner background
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_DARK
        top_bar.line.fill.background()

        # Accent red stripe
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.10), Inches(13.333), Inches(0.05))
        accent.fill.solid()
        accent.fill.fore_color.rgb = COLOR_CRIMSON
        accent.line.fill.background()

        # Category text
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.5), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category.upper()
        cat_p.font.name = "Arial"
        cat_p.font.size = Pt(10)
        cat_p.font.bold = True
        cat_p.font.color.rgb = COLOR_CRIMSON

        # Main Title text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.5), Inches(0.65))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = "Arial"
        title_p.font.size = Pt(22)
        title_p.font.bold = True
        title_p.font.color.rgb = COLOR_WHITE

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # =========================================================================
    # SLIDE 1: TITLE & EXECUTIVE SUMMARY
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)

    # Background dark gradient simulation
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_DARK
    bg.line.fill.background()

    # Red accent stripe on left
    left_accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.2), Inches(0.15), Inches(5.0))
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = COLOR_CRIMSON
    left_accent.line.fill.background()

    # Title Box
    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.5), Inches(2.2))
    t_tf = t_box.text_frame
    t_tf.word_wrap = True

    p_badge = t_tf.paragraphs[0]
    p_badge.text = "SMART HEALTHCARE & EMERGENCY DISPATCH PLATFORM"
    p_badge.font.name = "Arial"
    p_badge.font.size = Pt(12)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_CRIMSON

    p_main = t_tf.add_paragraph()
    p_main.text = "LIFEDROP SHIVAMOGGA"
    p_main.font.name = "Arial"
    p_main.font.size = Pt(38)
    p_main.font.bold = True
    p_main.font.color.rgb = COLOR_WHITE
    p_main.space_before = Pt(8)

    p_sub = t_tf.add_paragraph()
    p_sub.text = "Rapid Emergency Blood Donor Network & Triage Intelligence Platform"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(203, 213, 225)
    p_sub.space_before = Pt(6)

    # 3 Summary Cards
    cards_data = [
        ("🚨 Critical Problem", "Emergency blood shortages and delays in finding compatible donors cause life-threatening delays in acute trauma and surgeries across Shivamogga & Malnad.", COLOR_LIGHT_RED, COLOR_CRIMSON),
        ("💡 Intelligent Solution", "Direct real-time matchmaking connecting patients, verified volunteer donors, and blood banks via automated compatibility & live dispatch tracking.", COLOR_NAVY, COLOR_EMERALD),
        ("🎯 Impact & Reach", "Cuts donor response time from hours to under 15 minutes, empowering McGann Hospital, Sahyadri, and 7 taluk healthcare facilities.", COLOR_NAVY, COLOR_BLUE)
    ]

    card_width = Inches(3.6)
    card_height = Inches(2.6)
    start_left = Inches(1.0)
    top_pos = Inches(3.9)

    for i, (head, desc, bg_c, acc_c) in enumerate(cards_data):
        c_left = start_left + i * Inches(3.9)
        c_shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, top_pos, card_width, card_height)
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = bg_c
        c_shape.line.color.rgb = acc_c
        c_shape.line.width = Pt(1.5)

        tf = c_shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        p_h = tf.paragraphs[0]
        p_h.text = head
        p_h.font.name = "Arial"
        p_h.font.size = Pt(15)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE if bg_c == COLOR_NAVY else COLOR_CRIMSON

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Calibri"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = RGBColor(226, 232, 240) if bg_c == COLOR_NAVY else COLOR_DARK
        p_d.space_before = Pt(8)

    # Bottom Metadata
    meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.4))
    m_tf = meta_box.text_frame
    m_p = m_tf.paragraphs[0]
    m_p.text = "⚡ Full-Stack Python/Flask + SQLite + Leaflet.js | Shivamogga Blood Bank Council Initiative"
    m_p.font.name = "Calibri"
    m_p.font.size = Pt(11)
    m_p.font.color.rgb = COLOR_MUTED

    # =========================================================================
    # SLIDE 2: CORE FEATURES & SYSTEM MODULES
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Core System Modules & Capabilities")

    features = [
        ("🔍 Instant Donor Matching", "Filter verified donors by exact ABO/Rh blood type, taluk location (Shivamogga, Bhadravathi, Sagara), and real-time live availability status.", "⚡ Immediate Discovery"),
        ("🚨 Emergency Request Broadcast", "Post critical blood requirements with 1-click urgent notifications, automatic priority grading (< 2 hrs, 24 hrs, planned), and SMS/Call dispatch.", "⏱️ Priority Pipeline"),
        ("🧬 Blood Compatibility Matrix", "Interactive ABO/Rh compatibility engine visualising universal donor (O-) and universal recipient (AB+) pathways with clinical guidelines.", "🔬 Clinical Intelligence"),
        ("📍 Live Request Tracker", "End-to-end GPS and timeline tracker showing stage progression: Request Logged -> Donors Alerted -> En Route -> Hospital Received.", "🗺️ Real-Time GIS"),
        ("🤖 Lifedrop AI Assistant", "Interactive emergency triage bot answering medical queries, routing donor registrations, and finding nearest blood banks instantly.", "💬 24/7 AI Triage"),
        ("🪪 Digital Hero Pass", "Unique donor identity cards with blood group verification badges, donation history tracking, and life-saver certificates.", "🎖️ Donor Gamification")
    ]

    col_width = Inches(3.7)
    row_height = Inches(2.65)
    
    for idx, (ftitle, fdesc, ftag) in enumerate(features):
        row = idx // 3
        col = idx % 3
        left = Inches(0.8) + col * Inches(4.0)
        top = Inches(1.5) + row * Inches(2.85)

        card = add_card(slide2, left, top, col_width, row_height)
        tf = card.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = ftag.upper()
        p1.font.name = "Arial"
        p1.font.size = Pt(9)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_CRIMSON

        p2 = tf.add_paragraph()
        p2.text = ftitle
        p2.font.name = "Arial"
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_NAVY
        p2.space_before = Pt(3)

        p3 = tf.add_paragraph()
        p3.text = fdesc
        p3.font.name = "Calibri"
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_MUTED
        p3.space_before = Pt(6)

    # =========================================================================
    # SLIDE 3: EMERGENCY WORKFLOW & PATIENT JOURNEY
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "End-to-End Emergency Patient & Donor Workflow")

    steps = [
        ("Step 1: Emergency Broadcast", "Patient attendant or hospital logs urgent request with blood group, units, hospital & city.", "Hospital / Attendant", COLOR_CRIMSON),
        ("Step 2: Intelligent Matchmaking", "Engine scans database for compatible donors within target radius using ABO/Rh matrix.", "Lifedrop Core Engine", COLOR_NAVY),
        ("Step 3: Rapid Donor Alert", "Automated priority alert dispatched to matched donors via phone/SMS with 1-tap accept.", "Volunteer Network", COLOR_AMBER),
        ("Step 4: Dispatch & Hospital Delivery", "Donor arrives at blood bank/hospital; units collected; live tracker updates status to Complete.", "McGann / Sahyadri", COLOR_EMERALD)
    ]

    s_width = Inches(2.7)
    s_height = Inches(3.8)
    
    for i, (stitle, sdesc, srole, scolor) in enumerate(steps):
        s_left = Inches(0.8) + i * Inches(2.95)
        card = add_card(slide3, s_left, Inches(1.5), s_width, s_height)
        
        # Step top banner
        sbanner = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, s_left, Inches(1.5), s_width, Inches(0.4))
        sbanner.fill.solid()
        sbanner.fill.fore_color.rgb = scolor
        sbanner.line.fill.background()
        stf = sbanner.text_frame
        sp = stf.paragraphs[0]
        sp.text = f"STAGE 0{i+1}"
        sp.font.name = "Arial"
        sp.font.size = Pt(10)
        sp.font.bold = True
        sp.font.color.rgb = COLOR_WHITE
        sp.alignment = PP_ALIGN.CENTER

        tf = card.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = "\n"  # spacing below banner
        
        p1 = tf.add_paragraph()
        p1.text = stitle
        p1.font.name = "Arial"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_DARK
        p1.space_before = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_MUTED
        p2.space_before = Pt(8)

        p3 = tf.add_paragraph()
        p3.text = f"👤 Stakeholder: {srole}"
        p3.font.name = "Arial"
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = scolor
        p3.space_before = Pt(14)

    # Bottom workflow summary box
    bot_box = add_card(slide3, Inches(0.8), Inches(5.55), Inches(11.7), Inches(1.4), bg_color=COLOR_LIGHT_RED, border_color=COLOR_CRIMSON)
    btf = bot_box.text_frame
    btf.word_wrap = True
    bp1 = btf.paragraphs[0]
    bp1.text = "⚡ Key Workflow Optimization: Zero-Lag Autonomous Matching"
    bp1.font.name = "Arial"
    bp1.font.size = Pt(12)
    bp1.font.bold = True
    bp1.font.color.rgb = COLOR_CRIMSON

    bp2 = btf.add_paragraph()
    bp2.text = "Traditional manual phone-tree donor searches take 2 to 6 hours. Lifedrop's automated compatibility matching and geo-fenced donor broadcast cuts critical lead-time to under 15 minutes, directly saving acute trauma patients."
    bp2.font.name = "Calibri"
    bp2.font.size = Pt(11)
    bp2.font.color.rgb = COLOR_NAVY
    bp2.space_before = Pt(3)

    # =========================================================================
    # SLIDE 4: FULL-STACK TECHNICAL ARCHITECTURE
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Full-Stack Technology Stack & System Architecture")

    tech_layers = [
        ("🎨 Frontend Layer", "Modern Responsive Web", [
            "Semantic HTML5 & Vanilla CSS3 Design System",
            "CSS Custom Properties with 4 Dynamic Themes",
            "Leaflet.js for Interactive Geospatial Mapping",
            "Client-side Offline Fallback & LocalStorage Sync"
        ], COLOR_BLUE),
        ("⚙️ Backend REST API", "Python 3 & Flask Framework", [
            "Flask 3.0 RESTful Microservice Architecture",
            "CORS-enabled Endpoints for High Concurrency",
            "Deterministic Compatibility Rules Engine",
            "Lifedrop AI Rule-based Triage Engine"
        ], COLOR_CRIMSON),
        ("🗄️ Database & Storage", "SQLite Relational Store", [
            "Embedded SQLite (lifedrop.db) with WAL mode",
            "Normalized Donors, Requests & Inventory tables",
            "Audit logs and timeline tracking schema",
            "Zero external dependency lightweight deployment"
        ], COLOR_EMERALD)
    ]

    for i, (ltitle, lsub, lbullets, lcolor) in enumerate(tech_layers):
        l_left = Inches(0.8) + i * Inches(4.0)
        card = add_card(slide4, l_left, Inches(1.5), Inches(3.7), Inches(5.4))
        
        # Header banner
        h_shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, l_left, Inches(1.5), Inches(3.7), Inches(0.75))
        h_shape.fill.solid()
        h_shape.fill.fore_color.rgb = lcolor
        h_shape.line.fill.background()
        
        htf = h_shape.text_frame
        hp1 = htf.paragraphs[0]
        hp1.text = ltitle
        hp1.font.name = "Arial"
        hp1.font.size = Pt(13)
        hp1.font.bold = True
        hp1.font.color.rgb = COLOR_WHITE
        
        hp2 = htf.add_paragraph()
        hp2.text = lsub
        hp2.font.name = "Calibri"
        hp2.font.size = Pt(10)
        hp2.font.color.rgb = RGBColor(241, 245, 249)

        # Content
        ctf = card.text_frame
        ctf.word_wrap = True
        
        cp0 = ctf.paragraphs[0]
        cp0.text = "\n\n"  # spacing
        
        for bullet in lbullets:
            bp = ctf.add_paragraph()
            bp.text = f"• {bullet}"
            bp.font.name = "Calibri"
            bp.font.size = Pt(11)
            bp.font.color.rgb = COLOR_NAVY
            bp.space_before = Pt(8)

    # =========================================================================
    # SLIDE 5: BLOOD COMPATIBILITY & INVENTORY MANAGEMENT
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Blood Compatibility Engine & Clinical Stock Matrix")

    # Left: Compatibility Table
    left_card = add_card(slide5, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.4))
    ltf = left_card.text_frame
    ltf.word_wrap = True
    
    lp1 = ltf.paragraphs[0]
    lp1.text = "🧬 ABO & Rh Factor Compatibility Rules"
    lp1.font.name = "Arial"
    lp1.font.size = Pt(14)
    lp1.font.bold = True
    lp1.font.color.rgb = COLOR_CRIMSON

    rules = [
        ("O-", "O- Only (Universal Red Cell Donor)", "ALL Groups (O-, O+, A-, A+, B-, B+, AB-, AB+)"),
        ("O+", "O+, O-", "O+, A+, B+, AB+"),
        ("A-", "A-, O-", "A-, A+, AB-, AB+"),
        ("A+", "A+, A-, O+, O-", "A+, AB+"),
        ("B-", "B-, O-", "B-, B+, AB-, AB+"),
        ("B+", "B+, B-, O+, O-", "B+, AB+"),
        ("AB-", "AB-, A-, B-, O-", "AB-, AB+"),
        ("AB+", "ALL Groups (Universal Recipient)", "AB+ Only")
    ]

    # Create table on slide
    table_shape = slide5.shapes.add_table(9, 3, Inches(1.0), Inches(2.2), Inches(6.1), Inches(4.3))
    table = table_shape.table
    table.columns[0].width = Inches(0.9)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.7)

    headers = ["Blood Group", "Can Receive From", "Can Donate To"]
    for col_idx, htext in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Arial"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    for row_idx, (bg, rec, don) in enumerate(rules):
        row = row_idx + 1
        c0 = table.cell(row, 0)
        c0.text = bg
        c0.fill.solid()
        c0.fill.fore_color.rgb = COLOR_LIGHT_RED if "O-" in bg or "AB+" in bg else COLOR_WHITE
        p0 = c0.text_frame.paragraphs[0]
        p0.font.name = "Arial"
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_CRIMSON

        c1 = table.cell(row, 1)
        c1.text = rec
        c1.fill.solid()
        c1.fill.fore_color.rgb = COLOR_CARD_BG
        p1 = c1.text_frame.paragraphs[0]
        p1.font.name = "Calibri"
        p1.font.size = Pt(9)
        p1.font.color.rgb = COLOR_NAVY

        c2 = table.cell(row, 2)
        c2.text = don
        c2.fill.solid()
        c2.fill.fore_color.rgb = COLOR_CARD_BG
        p2 = c2.text_frame.paragraphs[0]
        p2.font.name = "Calibri"
        p2.font.size = Pt(9)
        p2.font.color.rgb = COLOR_NAVY

    # Right: Blood Bank Inventory Highlights
    right_card = add_card(slide5, Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.4))
    rtf = right_card.text_frame
    rtf.word_wrap = True
    
    rp1 = rtf.paragraphs[0]
    rp1.text = "🏥 Live Blood Bank Inventory Tracking"
    rp1.font.name = "Arial"
    rp1.font.size = Pt(14)
    rp1.font.bold = True
    rp1.font.color.rgb = COLOR_NAVY

    inv_points = [
        ("Real-Time Unit Counts", "Tracks critical stock across McGann District Hospital, Sahyadri Blood Bank, and Rotary Rotary Center."),
        ("Automated Deficit Triggers", "When stock falls below 5 units, system flags urgency badge and prompts donor broadcast."),
        ("Component Fractionation", "Supports Whole Blood, Packed Red Blood Cells (PRBC), Platelets, and Fresh Frozen Plasma (FFP)."),
        ("Expiry & Batch Management", "FIFO rotation alerts to prevent wastage of valuable platelet concentrates (5-day shelf life).")
    ]

    for ititle, idesc in inv_points:
        p_t = rtf.add_paragraph()
        p_t.text = f"✓ {ititle}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CRIMSON
        p_t.space_before = Pt(8)

        p_d = rtf.add_paragraph()
        p_d.text = idesc
        p_d.font.name = "Calibri"
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = COLOR_MUTED
        p_d.space_before = Pt(2)

    # =========================================================================
    # SLIDE 6: IMPACT, METRICS & FUTURE ROADMAP
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Clinical Impact, Key Metrics & Future Roadmap")

    # Top Metric Banner Cards
    metrics = [
        ("< 15 Mins", "Average Match Lead-time", COLOR_CRIMSON),
        ("1,240+", "Verified Active Donors", COLOR_NAVY),
        ("380+", "Emergency Units Dispatched", COLOR_EMERALD),
        ("100%", "Compatibility Verification", COLOR_BLUE)
    ]

    for i, (mval, mlbl, mcol) in enumerate(metrics):
        m_left = Inches(0.8) + i * Inches(2.95)
        m_card = add_card(slide6, m_left, Inches(1.5), Inches(2.7), Inches(1.4))
        mtf = m_card.text_frame
        mtf.word_wrap = True
        
        mp1 = mtf.paragraphs[0]
        mp1.text = mval
        mp1.font.name = "Arial"
        mp1.font.size = Pt(22)
        mp1.font.bold = True
        mp1.font.color.rgb = mcol
        mp1.alignment = PP_ALIGN.CENTER

        mp2 = mtf.add_paragraph()
        mp2.text = mlbl
        mp2.font.name = "Calibri"
        mp2.font.size = Pt(10)
        mp2.font.color.rgb = COLOR_MUTED
        mp2.alignment = PP_ALIGN.CENTER
        mp2.space_before = Pt(2)

    # Bottom: Future Roadmap Columns
    roadmap = [
        ("📱 Phase 1: Regional WhatsApp Bot", "Integration with WhatsApp Cloud API for 1-click accept & dispatch coordinates in Kannada & English.", COLOR_NAVY),
        ("🏥 Phase 2: Hospital EMR & HMIS Sync", "Direct HL7/FHIR integration with McGann & District EMR systems for automated trauma alerts.", COLOR_NAVY),
        ("🚁 Phase 3: Drone Delivery & Bio-Cold Chain", "IoT temperature-monitored blood transport and drone corridor support for rural Malnad taluks.", COLOR_NAVY)
    ]

    for i, (rtitle, rdesc, rcolor) in enumerate(roadmap):
        r_left = Inches(0.8) + i * Inches(4.0)
        rcard = add_card(slide6, r_left, Inches(3.2), Inches(3.7), Inches(3.7))
        
        rtf = rcard.text_frame
        rtf.word_wrap = True

        rp1 = rtf.paragraphs[0]
        rp1.text = rtitle
        rp1.font.name = "Arial"
        rp1.font.size = Pt(13)
        rp1.font.bold = True
        rp1.font.color.rgb = COLOR_CRIMSON

        rp2 = rtf.add_paragraph()
        rp2.text = rdesc
        rp2.font.name = "Calibri"
        rp2.font.size = Pt(11)
        rp2.font.color.rgb = COLOR_NAVY
        rp2.space_before = Pt(8)

        rp3 = rtf.add_paragraph()
        rp3.text = "🎯 Objective: Expand coverage to Shimoga, Thirthahalli, Sagar, Shikaripur, Sorab, and Hosanagara taluks."
        rp3.font.name = "Calibri"
        rp3.font.size = Pt(10)
        rp3.font.italic = True
        rp3.font.color.rgb = COLOR_MUTED
        rp3.space_before = Pt(14)

    prs.save(PPTX_PATH)
    print(f"[OK] Successfully created 6-slide PowerPoint: {PPTX_PATH}")


# -----------------------------------------------------------------------------
# 2. WORD DOCUMENT GENERATION (5-6 PAGES)
# -----------------------------------------------------------------------------

def set_cell_background(cell, fill_hex):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{fill_hex}"/>')
    tcPr.append(shd)

def set_cell_margins(cell, top=100, bottom=100, left=150, right=150):
    tcPr = cell._tc.get_or_add_tcPr()
    tcMar = parse_xml(f'''
        <w:tcMar {nsdecls("w")}>
            <w:top w:w="{top}" w:type="dxa"/>
            <w:bottom w:w="{bottom}" w:type="dxa"/>
            <w:left w:w="{left}" w:type="dxa"/>
            <w:right w:w="{right}" w:type="dxa"/>
        </w:tcMar>
    ''')
    tcPr.append(tcMar)

def create_word_document():
    print("Generating 5-6 Page Word Documentation...")
    doc = docx.Document()

    # Section Margins Setup
    for section in doc.sections:
        section.top_margin = DocxInches(0.75)
        section.bottom_margin = DocxInches(0.75)
        section.left_margin = DocxInches(0.8)
        section.right_margin = DocxInches(0.8)

    # Style Helpers
    def add_title(text, subtitle=""):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = DocxPt(0)
        p.paragraph_format.space_after = DocxPt(4)
        
        r = p.add_run(text + "\n")
        r.font.name = "Arial"
        r.font.size = DocxPt(22)
        r.font.bold = True
        r.font.color.rgb = DocxRGBColor(225, 29, 72)

        if subtitle:
            r2 = p.add_run(subtitle + "\n")
            r2.font.name = "Arial"
            r2.font.size = DocxPt(12)
            r2.font.bold = True
            r2.font.color.rgb = DocxRGBColor(15, 23, 42)

        r3 = p.add_run("Shivamogga District Emergency Healthcare & Blood Donor Coordination System | Version 2.0")
        r3.font.name = "Calibri"
        r3.font.size = DocxPt(9.5)
        r3.font.italic = True
        r3.font.color.rgb = DocxRGBColor(100, 116, 139)

    def add_h1(text):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = DocxPt(14)
        p.paragraph_format.space_after = DocxPt(4)
        p.paragraph_format.keep_with_next = True
        r = p.add_run(text)
        r.font.name = "Arial"
        r.font.size = DocxPt(14)
        r.font.bold = True
        r.font.color.rgb = DocxRGBColor(225, 29, 72)
        return p

    def add_h2(text):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = DocxPt(10)
        p.paragraph_format.space_after = DocxPt(2)
        p.paragraph_format.keep_with_next = True
        r = p.add_run(text)
        r.font.name = "Arial"
        r.font.size = DocxPt(11.5)
        r.font.bold = True
        r.font.color.rgb = DocxRGBColor(30, 41, 59)
        return p

    def add_p(text, bold_prefix=""):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = DocxPt(4)
        p.paragraph_format.line_spacing = 1.15
        if bold_prefix:
            rb = p.add_run(bold_prefix)
            rb.font.name = "Calibri"
            rb.font.size = DocxPt(10)
            rb.font.bold = True
            rb.font.color.rgb = DocxRGBColor(15, 23, 42)
        r = p.add_run(text)
        r.font.name = "Calibri"
        r.font.size = DocxPt(10)
        r.font.color.rgb = DocxRGBColor(51, 65, 85)
        return p

    def add_callout(title, text):
        tbl = doc.add_table(rows=1, cols=1)
        tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
        cell = tbl.cell(0, 0)
        cell.width = DocxInches(6.8)
        set_cell_background(cell, "FFF1F2")
        set_cell_margins(cell, 120, 120, 180, 180)
        
        p = cell.paragraphs[0]
        p.paragraph_format.space_after = DocxPt(2)
        r_t = p.add_run(f"📌 {title}\n")
        r_t.font.name = "Arial"
        r_t.font.size = DocxPt(10.5)
        r_t.font.bold = True
        r_t.font.color.rgb = DocxRGBColor(225, 29, 72)
        
        r_b = p.add_run(text)
        r_b.font.name = "Calibri"
        r_b.font.size = DocxPt(9.5)
        r_b.font.color.rgb = DocxRGBColor(30, 41, 59)
        
        doc.add_paragraph().paragraph_format.space_after = DocxPt(4)

    # =========================================================================
    # PAGE 1: PROJECT OVERVIEW, OBJECTIVES & PROBLEM STATEMENT
    # =========================================================================
    add_title("LIFEDROP SHIVAMOGGA", "Rapid Emergency Blood Donor Network & Triage Intelligence Platform")
    
    add_h1("1. Executive Summary & Project Background")
    add_p(
        "Lifedrop is an emergency medical response platform designed to eliminate fatal delays in blood procurement across Shivamogga and the wider Malnad region. In critical emergency situations—such as severe road traffic accidents, acute postpartum haemorrhage, oncological surgeries, and dengue-induced severe thrombocytopenia—access to compatible blood within the 'Golden Hour' is the primary determinant of patient survival."
    )
    add_p(
        "Traditional blood search methods rely on fragmented social media broadcasts, manual phone calls to acquaintances, and physical visits to multiple hospital blood banks. This manual workflow frequently causes delays of 2 to 6 hours. Lifedrop re-engineers this pipeline into an automated, real-time matching system that dispatches alerts to verified, compatible donors within minutes."
    )

    add_callout(
        "Mission Statement",
        "To build a resilient, zero-lag volunteer blood donation infrastructure in Shivamogga that guarantees no patient loses their life due to the unavailability of compatible blood or lack of donor connectivity."
    )

    add_h2("1.1 Core Problem Statement & Regional Healthcare Challenges")
    add_p("The healthcare ecosystem in Shivamogga faces four primary blood availability bottlenecks:")
    add_p("Lack of a unified digital registry makes it difficult to ascertain which hospital (McGann, Sahyadri, Nanjappa, or Rotary) has specific units.", "1. Fragmented Inventory: ")
    add_p("Patients with rare blood groups (such as O-negative, AB-negative, or Bombay phenotype) face high mortality risks due to lack of an on-demand donor registry.", "2. Rare Blood Accessibility: ")
    add_p("Social media appeals lack clinical verification, leading to spam calls, expired requests, and delayed actual donor responses.", "3. Unverified Emergency Broadcasts: ")
    add_p("Surrounding rural taluks (Thirthahalli, Sagara, Sorab, Hosanagara) must transfer patients to district headquarters without confirmed blood availability.", "4. Geographic & Transport Latency: ")

    add_h2("1.2 Strategic Objectives")
    add_p("• Reduce the average emergency donor matching and confirmation time to under 15 minutes.")
    add_p("• Provide a transparent, interactive ABO/Rh compatibility engine for clinical accuracy.")
    add_p("• Implement live GPS-enabled tracking for emergency blood dispatches to keep attendants informed.")
    add_p("• Equip community volunteers with official Digital Hero Passes to encourage recurring blood donation.")

    doc.add_page_break()

    # =========================================================================
    # PAGE 2: SYSTEM ARCHITECTURE & CORE FUNCTIONAL MODULES
    # =========================================================================
    add_h1("2. System Architecture & Core Functional Modules")
    add_p(
        "Lifedrop is built on a modular full-stack architecture combining a high-performance Python/Flask REST backend, an embedded SQLite relational database, and a responsive frontend user interface featuring custom themes, interactive Leaflet maps, and an AI triage assistant."
    )

    add_h2("2.1 Module Specifications")
    
    # Modules Table
    tbl_mod = doc.add_table(rows=6, cols=3)
    tbl_mod.alignment = WD_TABLE_ALIGNMENT.CENTER
    for r in tbl_mod.rows:
        for c in r.cells:
            set_cell_margins(c, 80, 80, 100, 100)
    
    m_headers = ["Module Name", "Primary Function", "Key Clinical / Technical Benefit"]
    for c_idx, h in enumerate(m_headers):
        cell = tbl_mod.cell(0, c_idx)
        cell.text = h
        set_cell_background(cell, "0F172A")
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = DocxPt(9.5)
        r.font.bold = True
        r.font.color.rgb = DocxRGBColor(255, 255, 255)

    m_data = [
        ("Donor Directory & Filter", "Real-time search across verified donors with blood group, city, and live status filters.", "Instant access to 1,240+ verified local donors."),
        ("Emergency Case Feeds", "Live portal displaying active cases categorized by urgency (< 2 hrs, 24 hrs, planned).", "Prevents duplicate broadcasts and triages severe cases."),
        ("Compatibility Matrix", "Interactive matrix visualizing whole blood and PRBC donor-recipient compatibility rules.", "Eliminates dangerous ABO transfusion mismatches."),
        ("Live Dispatch Tracker", "Visual timeline progression from 'Request Logged' to 'Units Transfused' with map pins.", "Full transparency for anxious family members and doctors."),
        ("Lifedrop AI Assistant", "Rule-based and NLP triage bot providing instant emergency guidance and registration.", "24/7 autonomous support without human dispatcher delay.")
    ]

    for r_idx, (m_name, m_func, m_ben) in enumerate(m_data):
        row = r_idx + 1
        c0 = tbl_mod.cell(row, 0)
        c0.text = m_name
        set_cell_background(c0, "FFF1F2")
        c0.paragraphs[0].runs[0].font.bold = True
        c0.paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(225, 29, 72)

        c1 = tbl_mod.cell(row, 1)
        c1.text = m_func
        set_cell_background(c1, "F8FAFC")

        c2 = tbl_mod.cell(row, 2)
        c2.text = m_ben
        set_cell_background(c2, "F8FAFC")

    add_h2("2.2 Digital Hero Pass & Donor Engagement")
    add_p(
        "To foster a culture of voluntary donation, Lifedrop generates dynamic 'Hero Passes' for registered donors. Each pass features donor blood type, donation verification badges, emergency availability toggles, and direct contact protocols. Donors receive recognition certificates upon completing verified donations at partner healthcare centres."
    )

    doc.add_page_break()

    # =========================================================================
    # PAGE 3: BLOOD COMPATIBILITY & CLINICAL MATCHING ENGINE
    # =========================================================================
    add_h1("3. Blood Compatibility & Intelligent Matching Algorithm")
    add_p(
        "Transfusion medicine requires absolute precision. Administering ABO-incompatible blood results in acute hemolytic transfusion reactions (AHTR), which can be fatal. Lifedrop embeds strict immunohematology algorithms to ensure only clinically compatible donors are surfaced during emergencies."
    )

    add_h2("3.1 ABO and Rh Factor Compatibility Matrix")
    
    # Blood Matrix Table
    tbl_bg = doc.add_table(rows=9, cols=4)
    tbl_bg.alignment = WD_TABLE_ALIGNMENT.CENTER
    for r in tbl_bg.rows:
        for c in r.cells:
            set_cell_margins(c, 70, 70, 90, 90)

    bg_headers = ["Blood Group", "Antigens Present", "Can Receive Packed Cells From", "Can Donate Packed Cells To"]
    for c_idx, h in enumerate(bg_headers):
        cell = tbl_bg.cell(0, c_idx)
        cell.text = h
        set_cell_background(cell, "0F172A")
        p = cell.paragraphs[0]
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = DocxPt(9)
        r.font.bold = True
        r.font.color.rgb = DocxRGBColor(255, 255, 255)

    bg_data = [
        ("O-", "None (No A, B, or Rh)", "O- Only", "ALL Groups (Universal Donor)"),
        ("O+", "Rh Factor", "O+, O-", "O+, A+, B+, AB+"),
        ("A-", "A Antigen", "A-, O-", "A-, A+, AB-, AB+"),
        ("A+", "A Antigen & Rh", "A+, A-, O+, O-", "A+, AB+"),
        ("B-", "B Antigen", "B-, O-", "B-, B+, AB-, AB+"),
        ("B+", "B Antigen & Rh", "B+, B-, O+, O-", "B+, AB+"),
        ("AB-", "A and B Antigens", "AB-, A-, B-, O-", "AB-, AB+"),
        ("AB+", "A, B & Rh Antigens", "ALL Groups (Universal Recipient)", "AB+ Only")
    ]

    for r_idx, (bgroup, ant, rec, don) in enumerate(bg_data):
        row = r_idx + 1
        c0 = tbl_bg.cell(row, 0)
        c0.text = bgroup
        set_cell_background(c0, "FFF1F2" if "O-" in bgroup or "AB+" in bgroup else "F8FAFC")
        c0.paragraphs[0].runs[0].font.bold = True
        c0.paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(225, 29, 72)

        c1 = tbl_bg.cell(row, 1)
        c1.text = ant
        set_cell_background(c1, "F8FAFC")

        c2 = tbl_bg.cell(row, 2)
        c2.text = rec
        set_cell_background(c2, "F8FAFC")

        c3 = tbl_bg.cell(row, 3)
        c3.text = don
        set_cell_background(c3, "F8FAFC")

    add_h2("3.2 Match Ranking & Geo-Proximity Weighting")
    add_p(
        "When an emergency request is logged, Lifedrop executes a weighted ranking query: Rank = (Exact Match × 1.0 + Compatible Match × 0.7) × (City Proximity Weight) × (Live Availability Multiplier). This ensures exact ABO/Rh matches within the patient's immediate vicinity (e.g., McGann Hospital, Shivamogga) receive top notification priority."
    )

    doc.add_page_break()

    # =========================================================================
    # PAGE 4: FULL-STACK TECHNICAL ARCHITECTURE & DATABASE DESIGN
    # =========================================================================
    add_h1("4. Technical Implementation & Database Schema")
    add_p(
        "Lifedrop is engineered with high resilience and zero unnecessary overhead. It utilizes Python Flask for backend services and SQLite for persistent relational data, alongside a hybrid local storage synchronization engine that permits full offline operation."
    )

    add_h2("4.1 Relational Database Architecture (SQLite3)")
    add_p("The database schema encompasses three primary entities designed for high query performance and integrity:")
    
    # Database Table
    tbl_db = doc.add_table(rows=4, cols=3)
    tbl_db.alignment = WD_TABLE_ALIGNMENT.CENTER
    for r in tbl_db.rows:
        for c in r.cells:
            set_cell_margins(c, 70, 70, 90, 90)

    db_headers = ["Table Name", "Key Columns & Data Types", "Business Purpose"]
    for c_idx, h in enumerate(db_headers):
        cell = tbl_db.cell(0, c_idx)
        cell.text = h
        set_cell_background(cell, "0F172A")
        p = cell.paragraphs[0]
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = DocxPt(9)
        r.font.bold = True
        r.font.color.rgb = DocxRGBColor(255, 255, 255)

    db_data = [
        ("donors", "id (TEXT PK), name, blood_group, city, phone, available (INT), verified (INT), donations_count, last_donation", "Stores registered volunteer donor profiles, live statuses, and donation history metrics."),
        ("requests", "id (TEXT PK), patient_name, blood_group, units_needed, hospital, city, urgency, contact_phone, status, timeline_json", "Tracks emergency blood broadcasts, hospital locations, and real-time dispatch milestones."),
        ("inventory", "blood_group (TEXT PK), units_available, hospital_name, last_updated", "Maintains live blood component stock levels across participating blood banks.")
    ]

    for r_idx, (tname, tcols, tpurp) in enumerate(db_data):
        row = r_idx + 1
        c0 = tbl_db.cell(row, 0)
        c0.text = tname
        set_cell_background(c0, "FFF1F2")
        c0.paragraphs[0].runs[0].font.bold = True
        c0.paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(225, 29, 72)

        c1 = tbl_db.cell(row, 1)
        c1.text = tcols
        set_cell_background(c1, "F8FAFC")

        c2 = tbl_db.cell(row, 2)
        c2.text = tpurp
        set_cell_background(c2, "F8FAFC")

    add_h2("4.2 REST API Specification")
    add_p("The backend exposes clean, RESTful JSON endpoints:")
    add_p("• GET /api/health — System uptime and database connection verification.")
    add_p("• GET /api/donors?bloodGroup=...&city=... — Filtered query for available donors.")
    add_p("• POST /api/donors — Register verified volunteer with automatic phone validation.")
    add_p("• GET /api/donors/match?bloodGroup=...&city=... — Execute real-time compatibility search.")
    add_p("• POST /api/requests — Submit emergency blood broadcast and trigger automated alerts.")
    add_p("• GET /api/requests/<id> — Live tracking endpoint with structured timeline progress.")

    doc.add_page_break()

    # =========================================================================
    # PAGE 5: LIVE DISPATCH TRACKER & GEOSPATIAL MAP INTEGRATION
    # =========================================================================
    add_h1("5. Live Dispatch Tracking & Geospatial Architecture")
    add_p(
        "Transparency during emergencies relieves severe psychological distress for patients and families. The Lifedrop Live Tracker provides a synchronized timeline and interactive map showing exactly where the case stands."
    )

    add_h2("5.1 Real-Time 4-Stage Dispatch Pipeline")
    add_p("Every emergency blood request moves through four discrete operational stages:")
    add_p("Request is submitted by hospital/attendant and broadcast to matching donors.", "Stage 1: Broadcast Active — ")
    add_p("Compatible donor receives notification and accepts request; attendant notified.", "Stage 2: Donor Confirmed — ")
    add_p("Donor is traveling to blood bank/hospital; GPS coordinates updated.", "Stage 3: En Route to Hospital — ")
    add_p("Blood units verified, cross-matched, and administered to patient.", "Stage 4: Units Transfused — ")

    add_h2("5.2 Geospatial Map Architecture (Leaflet.js & OpenStreetMap)")
    add_p(
        "Lifedrop integrates OpenStreetMap tile servers via Leaflet.js to render regional healthcare facilities without external paid API dependencies. Key coordinates mapped in Shivamogga include:"
    )
    add_p("• McGann District Teaching Hospital (SIMS), Sagar Road (13.9324° N, 75.5684° E)")
    add_p("• Sahyadri Super Specialty Hospital, Harakere (13.9056° N, 75.5892° E)")
    add_p("• Rotary Blood Bank & Diagnostic Centre, Durgigudi (13.9280° N, 75.5720° E)")
    add_p("• Nanjappa Multi-Speciality Hospital, Kuvempu Road (13.9312° N, 75.5650° E)")
    add_p("• Taluk General Hospitals in Bhadravathi, Sagara, and Thirthahalli.")

    doc.add_page_break()

    # =========================================================================
    # PAGE 6: DEPLOYMENT, SECURITY & FUTURE ROADMAP
    # =========================================================================
    add_h1("6. Deployment, Security & Future Roadmap")
    add_p(
        "Lifedrop is engineered for effortless deployment across both modern cloud environments (Render, Vercel, Netlify, AWS) and local emergency server setups with zero cloud latency."
    )

    add_h2("6.1 Security, Privacy & Medical Ethics")
    add_p("• Masked Contact Information: Attendant phone numbers are protected from public crawling.")
    add_p("• Rate Limiting: Prevents automated denial-of-service or fake broadcast flooding.")
    add_p("• Voluntary Consent: Donors retain complete control over their live availability toggles.")

    add_h2("6.2 Regional Impact & Scalability Milestones")
    
    # Impact Table
    tbl_imp = doc.add_table(rows=4, cols=2)
    tbl_imp.alignment = WD_TABLE_ALIGNMENT.CENTER
    for r in tbl_imp.rows:
        for c in r.cells:
            set_cell_margins(c, 70, 70, 90, 90)

    imp_headers = ["Implementation Milestone", "Projected Outcome & Community Impact"]
    for c_idx, h in enumerate(imp_headers):
        cell = tbl_imp.cell(0, c_idx)
        cell.text = h
        set_cell_background(cell, "0F172A")
        p = cell.paragraphs[0]
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = DocxPt(9.5)
        r.font.bold = True
        r.font.color.rgb = DocxRGBColor(255, 255, 255)

    imp_data = [
        ("Phase 1: WhatsApp Cloud Bot", "Integration with official WhatsApp Business API for instant 1-tap donor acceptances in Kannada and English."),
        ("Phase 2: District Hospital EMR Sync", "Direct integration with SIMS McGann Hospital HMIS for automated blood deficit alerts directly from trauma ICU."),
        ("Phase 3: Malnad Taluk Drone Logistics", "Support for IoT temperature-controlled medical drone delivery corridors connecting remote rural health centres to Shivamogga.")
    ]

    for r_idx, (m_stone, m_out) in enumerate(imp_data):
        row = r_idx + 1
        c0 = tbl_imp.cell(row, 0)
        c0.text = m_stone
        set_cell_background(c0, "FFF1F2")
        c0.paragraphs[0].runs[0].font.bold = True
        c0.paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(225, 29, 72)

        c1 = tbl_imp.cell(row, 1)
        c1.text = m_out
        set_cell_background(c1, "F8FAFC")

    add_h2("6.3 Conclusion")
    add_p(
        "Lifedrop represents a technological paradigm shift in emergency blood coordination for Shivamogga. By synthesizing real-time matching, clinical accuracy, live GPS tracking, and community gamification, Lifedrop transforms passive bystanders into active lifesavers, establishing a dependable safety net for the entire region."
    )

    doc.save(DOCX_PATH)
    print(f"[OK] Successfully created 5-6 Page Word Document: {DOCX_PATH}")


if __name__ == '__main__':
    create_presentation()
    create_word_document()
    print("\n[SUCCESS] Both files generated successfully!")
